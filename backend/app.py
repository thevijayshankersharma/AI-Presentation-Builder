import os
import logging
from flask import Flask, send_file, request, jsonify
from flask_cors import CORS
from pptx import Presentation
from pptx.util import Pt, Inches
from pptx.dml.color import RGBColor
from pptx.chart.data import CategoryChartData
from pptx.enum.chart import XL_CHART_TYPE
from pptx.oxml.xmlchemy import OxmlElement
from pptx.oxml.ns import qn
import google.generativeai as genai
import re
import requests
from io import BytesIO
import base64
import pandas as pd
from dotenv import load_dotenv

# Load environment variables
load_dotenv()

# Configure logging
logging.basicConfig(level=logging.DEBUG)

# Initialize Flask app and enable CORS
app = Flask(__name__)
CORS(app)

# Configure Gemini API
genai_api_key = os.environ.get('GENAI_API_KEY')
if not genai_api_key:
    app.logger.error("GENAI_API_KEY is not set in the environment!")
    raise ValueError("GENAI_API_KEY is required")
genai.configure(api_key=genai_api_key)

# Enhanced theme options with gradients and shadows
THEMES = {
    "corporate": {
        "gradient_start": RGBColor(0, 51, 102),  # Dark blue
        "gradient_end": RGBColor(173, 216, 230),  # Light blue
        "title_color": RGBColor(255, 255, 255),   # White
        "text_color": RGBColor(255, 255, 255),    # White
        "font_name": "Arial",
        "shadow": True,
    },
    "creative": {
        "gradient_start": RGBColor(147, 112, 219),  # Purple
        "gradient_end": RGBColor(255, 182, 193),    # Light pink
        "title_color": RGBColor(255, 255, 255),     # White
        "text_color": RGBColor(240, 240, 240),      # Off-white
        "font_name": "Montserrat",
        "shadow": True,
    },
    "minimal": {
        "gradient_start": RGBColor(245, 245, 245),  # Light gray
        "gradient_end": RGBColor(255, 255, 255),    # White
        "title_color": RGBColor(33, 33, 33),        # Dark gray
        "text_color": RGBColor(66, 66, 66),         # Medium gray
        "font_name": "Helvetica",
        "shadow": False,
    },
    "bold": {
        "gradient_start": RGBColor(255, 69, 0),     # Red-orange
        "gradient_end": RGBColor(255, 215, 0),      # Gold
        "title_color": RGBColor(255, 255, 255),     # White
        "text_color": RGBColor(255, 255, 255),      # White
        "font_name": "Impact",
        "shadow": True,
    }
}

# Chart type mapping
CHART_TYPES = {
    "bar": XL_CHART_TYPE.COLUMN_CLUSTERED,
    "line": XL_CHART_TYPE.LINE,
    "pie": XL_CHART_TYPE.PIE,
    "scatter": XL_CHART_TYPE.XY_SCATTER,
}

# Presentation styling constants
TITLE_FONT_SIZE = Pt(36)  # Larger for impact
CONTENT_FONT_SIZE = Pt(20)  # Slightly bigger for readability

# Function to apply gradient background
def apply_gradient(slide, start_color, end_color):
    background = slide.background
    fill = background.fill
    fill.gradient()
    fill.gradient_stops[0].color.rgb = start_color
    fill.gradient_stops[1].color.rgb = end_color
    fill.gradient_angle = 45  # Diagonal gradient

# Function to add shadow to text
def add_text_shadow(text_frame):
    for paragraph in text_frame.paragraphs:
        for run in paragraph.runs:
            add_shadow_to_run(run)

def add_shadow_to_run(run):
    rPr = run._r.get_or_add_rPr()
    effectLst = rPr.find(qn('a:effectLst'))
    if effectLst is None:
        effectLst = OxmlElement('a:effectLst')
        rPr.append(effectLst)
    
    # Check if outerShdw already exists
    outerShdw = effectLst.find(qn('a:outerShdw'))
    if outerShdw is None:
        outerShdw = OxmlElement('a:outerShdw')
        outerShdw.set('blurRad', '40000')  # Blur radius (~4pt)
        outerShdw.set('dist', '20000')     # Distance (~2pt)
        outerShdw.set('dir', '5400000')    # Direction (45 degrees)
        outerShdw.set('rotWithShape', '0')
        
        srgbClr = OxmlElement('a:srgbClr')
        srgbClr.set('val', '323232')  # Dark gray (RGB: 50, 50, 50)
        
        alpha = OxmlElement('a:alpha')
        alpha.set('val', '50000')  # 50% opacity
        srgbClr.append(alpha)
        
        outerShdw.append(srgbClr)
        effectLst.append(outerShdw)

def process_titles(text):
    try:
        lines = text.strip().split('\n')
        titles = [line.strip() for line in lines if line.strip()]
        return titles[:3]
    except Exception as e:
        app.logger.exception("Error processing titles")
        return [text]

def process_bullet_points(text):
    try:
        lines = text.split('\n')
        cleaned_lines = []
        for line in lines:
            line = line.strip()
            if line:
                line = re.sub(r'[\*\[\]]', '', line)
                line = re.sub(r'^[\d\.\-\s]+', '', line)
                cleaned_lines.append('- ' + line.strip())
        return '\n'.join(cleaned_lines)
    except Exception as e:
        app.logger.exception("Error processing bullet points")
        return text

def generate_slide_titles(content, language="en"):
    try:
        prompt = f"Generate exactly 3 concise slide titles for a presentation on '{content}' in {language}, no preamble or numbering."
        app.logger.debug(f"Generating titles with prompt: {prompt}")
        model = genai.GenerativeModel("gemini-1.5-flash")
        response = model.generate_content(prompt)
        titles_text = response.text.strip()
        app.logger.debug(f"Raw titles response: {titles_text}")
        titles = process_titles(titles_text)
        while len(titles) < 3:
            titles.append(f"{content} Overview {len(titles) + 1}")
        return titles[:3]
    except Exception as e:
        app.logger.exception(f"Failed to generate titles: {str(e)}")
        return [f"{content} Slide {i+1}" for i in range(3)]

def generate_slide_content(slide_title, has_image=True, summarize=False, language="en"):
    try:
        if summarize:
            prompt = f"Summarize content for '{slide_title}' into two concise paragraphs (max 30 words each) in {language}, no preamble or labels."
            max_tokens = 100
        elif has_image:
            prompt = f"Generate two concise paragraphs (max 50 words each) for '{slide_title}' in {language}, no preamble or labels."
            max_tokens = 150
        else:
            prompt = f"Generate six concise bullet points (max 25 words each) for '{slide_title}' in {language}. Use '-' as bullet marker, no numbering."
            max_tokens = 300
        app.logger.debug(f"Generating content with prompt: {prompt}")
        model = genai.GenerativeModel("gemini-1.5-flash")
        response = model.generate_content(prompt)
        content = response.text.strip()
        app.logger.debug(f"Raw content response: {content}")
        if not has_image and not summarize:
            content = process_bullet_points(content)
        return content
    except Exception as e:
        app.logger.exception(f"Failed to generate content for '{slide_title}': {str(e)}")
        return f"Content generation failed: {str(e)}"

def generate_image(prompt, language="en"):
    stability_api_key = os.environ.get('STABILITY_API_KEY')
    if not stability_api_key:
        app.logger.error("STABILITY_API_KEY is not set in the environment!")
        return None
    api_url = "https://api.stability.ai/v1/generation/stable-diffusion-xl-1024-v1-0/text-to-image"
    headers = {"Authorization": f"Bearer {stability_api_key}", "Content-Type": "application/json"}
    data = {
        "text_prompts": [{"text": f"{prompt}, professional high-quality illustration, styled for {language} audience"}],
        "cfg_scale": 7,
        "height": 1024,
        "width": 1024,
        "samples": 1,
        "steps": 30,
    }
    try:
        app.logger.debug(f"Requesting image with data: {data}")
        response = requests.post(api_url, headers=headers, json=data)
        response.raise_for_status()
        if "insufficient_balance" in response.text.lower():
            app.logger.error("Stability AI account has insufficient balance to generate image.")
            return None
        image_data = response.json()["artifacts"][0]["base64"]
        image_stream = BytesIO(base64.b64decode(image_data))
        app.logger.debug(f"Image generated successfully for prompt: {prompt}")
        return image_stream
    except Exception as e:
        app.logger.exception(f"Error generating image: {str(e)}")
        return None

def generate_chart(slide, csv_file, chart_type="bar"):
    try:
        df = pd.read_csv(csv_file)
        chart_data = CategoryChartData()
        chart_data.categories = df.iloc[:, 0].tolist()
        chart_data.add_series('Data', df.iloc[:, 1].tolist())
        chart_type_enum = CHART_TYPES.get(chart_type, XL_CHART_TYPE.COLUMN_CLUSTERED)
        chart = slide.shapes.add_chart(
            chart_type_enum,
            Inches(5.5), Inches(1.5),  # Right side position
            Inches(4.0), Inches(3.0),  # Adjusted size
            chart_data
        ).chart
        chart.has_title = True
        chart.chart_title.text_frame.text = "Data Overview"
        chart.chart_title.text_frame.paragraphs[0].font.size = Pt(14)
    except Exception as e:
        app.logger.exception("Error generating chart")

def create_presentation(topic, text_file, csv_file, theme="corporate", language="en", include_images=True, summarize=False, chart_type="bar"):
    try:
        prs = Presentation()
        selected_theme = THEMES.get(theme, THEMES["corporate"])

        content = topic
        if text_file:
            content = text_file.read().decode('utf-8')
        app.logger.debug(f"Content: {content[:100]}...")

        slide_titles = generate_slide_titles(content, language)
        
        # Title slide with enhanced design
        title_slide = prs.slides.add_slide(prs.slide_layouts[0])
        apply_gradient(title_slide, selected_theme["gradient_start"], selected_theme["gradient_end"])
        title_slide.shapes.title.text = slide_titles[0]
        tf = title_slide.shapes.title.text_frame
        tf.paragraphs[0].font.size = TITLE_FONT_SIZE
        tf.paragraphs[0].font.name = selected_theme["font_name"]
        tf.paragraphs[0].font.color.rgb = selected_theme["title_color"]
        tf.paragraphs[0].font.bold = True
        if selected_theme["shadow"]:
            add_text_shadow(tf)
        if len(title_slide.placeholders) > 1:
            subtitle = title_slide.placeholders[1]
            subtitle.text = "Powered by AI"
            subtitle.text_frame.paragraphs[0].font.size = Pt(24)
            subtitle.text_frame.paragraphs[0].font.name = selected_theme["font_name"]
            subtitle.text_frame.paragraphs[0].font.color.rgb = selected_theme["text_color"]
            if selected_theme["shadow"]:
                add_text_shadow(subtitle.text_frame)

        # Content slides
        for i, title in enumerate(slide_titles[1:], start=1):
            # Use layout 1 (title + content) for Slide 2 with image
            slide_layout = prs.slide_layouts[1] if (i == 1 and include_images) else prs.slide_layouts[5]
            slide = prs.slides.add_slide(slide_layout)
            apply_gradient(slide, selected_theme["gradient_start"], selected_theme["gradient_end"])

            slide.shapes.title.text = title
            tf = slide.shapes.title.text_frame
            tf.paragraphs[0].font.size = TITLE_FONT_SIZE
            tf.paragraphs[0].font.name = selected_theme["font_name"]
            tf.paragraphs[0].font.color.rgb = selected_theme["title_color"]
            tf.paragraphs[0].font.bold = True
            if selected_theme["shadow"]:
                add_text_shadow(tf)

            has_image = (i == 1 and include_images)
            has_chart = (csv_file and i == 2)
            content_width = Inches(5.0) if (has_image or has_chart) else Inches(9.0)  # Adjusted width
            content_text = generate_slide_content(title, has_image=has_image, summarize=summarize, language=language)

            if slide_layout == prs.slide_layouts[1] and has_image:
                # Use the content placeholder for text
                content_placeholder = slide.placeholders[1]
                content_placeholder.text = content_text
                tf = content_placeholder.text_frame
            else:
                # Manually add textbox with adjusted width
                textbox = slide.shapes.add_textbox(Inches(0.5), Inches(1.5), content_width, Inches(5.0))
                tf = textbox.text_frame
                tf.text = content_text

            tf.word_wrap = True
            for paragraph in tf.paragraphs:
                paragraph.font.size = CONTENT_FONT_SIZE
                paragraph.font.name = selected_theme["font_name"]
                paragraph.font.color.rgb = selected_theme["text_color"]
                paragraph.space_after = Pt(12)
                if selected_theme["shadow"]:
                    add_text_shadow(tf)

            if has_image:
                image_stream = generate_image(f"{title} related to {content_text}", language)
                if image_stream:
                    # Add image to the right side
                    slide.shapes.add_picture(image_stream, Inches(5.5), Inches(1.5), width=Inches(4.0))
                    app.logger.debug(f"Image added to slide: {title}")
                else:
                    app.logger.error(f"No image returned for slide: {title}")

            if has_chart:
                csv_file.seek(0)
                generate_chart(slide, csv_file, chart_type)

        output_dir = "generated_ppt"
        os.makedirs(output_dir, exist_ok=True)
        output_filepath = os.path.join(output_dir, f"{topic or 'presentation'}_presentation.pptx")
        prs.save(output_filepath)
        app.logger.debug(f"Saved to {output_filepath}")
        return output_filepath
    except Exception as e:
        app.logger.exception("Error in create_presentation")
        raise e

@app.route('/generate', methods=['POST'])
def generate():
    topic = request.form.get('topic')
    text_file = request.files.get('textFile')
    csv_file = request.files.get('csvFile')
    theme = request.form.get('theme', 'corporate')
    language = request.form.get('language', 'en')
    include_images = request.form.get('includeImages', 'true') == 'true'
    summarize = request.form.get('summarize', 'false') == 'true'
    chart_type = request.form.get('chartType', 'bar')
    export_format = request.form.get('exportFormat', 'pptx')

    if not topic and not text_file:
        return jsonify({"error": "Topic or text file required"}), 400

    try:
        app.logger.debug("Generating presentation")
        pptx_path = create_presentation(
            topic=topic,
            text_file=text_file,
            csv_file=csv_file,
            theme=theme,
            language=language,
            include_images=include_images,
            summarize=summarize,
            chart_type=chart_type
        )
        if export_format == 'pdf':
            return jsonify({"error": "PDF export not implemented"}), 501

        download_name = f"{topic or 'presentation'}.pptx"
        return send_file(pptx_path, as_attachment=True, download_name=download_name)
    except Exception as e:
        app.logger.exception("Error generating presentation")
        return jsonify({"error": str(e)}), 500

if __name__ == "__main__":
    app.run(port=5000, debug=True)